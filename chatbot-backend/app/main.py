# ============================
# main.py (Improved + Clean Version)
# ============================

import os
import re
import uuid
from datetime import datetime
from typing import Optional, List
import io

import httpx

from fastapi import FastAPI, UploadFile, File, Form, HTTPException
from fastapi.middleware.cors import CORSMiddleware
from pydantic import BaseModel, EmailStr
from pymongo import MongoClient
from groq import Groq
from dotenv import load_dotenv
import pdfplumber
import bcrypt
from docx import Document as DocxDocument
from pptx import Presentation

# ============================
# LOAD ENV
# ============================

from pathlib import Path
from dotenv import load_dotenv
import os

load_dotenv()

MONGO_URL = os.getenv("MONGO_URL")
GROQ_API_KEY = os.getenv("GROQ_API_KEY")
JINA_API_KEY = os.getenv("JINA_API_KEY")

if not MONGO_URL:
    raise Exception("❌ MONGO_URL missing in .env")

print("MongoDB:", "[OK] configured")
print("Groq:", "[OK] configured" if GROQ_API_KEY else "[MISSING]")
print("Jina:", "[OK] configured" if JINA_API_KEY else "[MISSING]")


# ============================
# FASTAPI SETUP
# ============================

app = FastAPI(title="RAG Chatbot Backend")

app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"],
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)


# ============================
# DATABASE CONNECTION
# ============================

try:
    client = MongoClient(MONGO_URL, serverSelectionTimeoutMS=5000, tls=True)
    client.server_info()
    db = client["rag_chatbot_db"]
    print("[OK] MongoDB connected")
except Exception as e:
    print("[ERROR] MongoDB failed:", e)
    db = None


users_collection = db["users"] if db is not None else []
documents_collection = db["documents"] if db is not None else []
responses_collection = db["responses"] if db is not None else []
chunks_collection = db["chunks"] if db is not None else []



# ============================
# GROQ CLIENT
# ============================

groq_client = Groq(api_key=GROQ_API_KEY) if GROQ_API_KEY else None


# ============================
# MODELS
# ============================

class UserRegister(BaseModel):
    username: str
    email: EmailStr
    password: str


class UserLogin(BaseModel):
    email: EmailStr
    password: str


class ChatRequest(BaseModel):
    user_id: str
    document_id: Optional[str] = None
    message: str


# ============================
# UTILS
# ============================

def hash_password(password: str):
    return bcrypt.hashpw(password.encode(), bcrypt.gensalt()).decode()


def verify_password(password: str, hashed: str):
    return bcrypt.checkpw(password.encode(), hashed.encode())


def serialize(doc):
    if not doc:
        return doc
    doc["_id"] = str(doc["_id"])
    for k, v in doc.items():
        if isinstance(v, datetime):
            doc[k] = v.isoformat()
    return doc


# ============================
# DOCUMENT EXTRACTION (PDF / DOCX / PPTX / TXT)
# ============================

SUPPORTED_EXTENSIONS = {".pdf", ".docx", ".pptx", ".txt"}


def _ext(filename: str) -> str:
    return os.path.splitext(filename or "")[1].lower()


def extract_text_from_pdf(file: UploadFile) -> str:
    text = ""
    with pdfplumber.open(file.file) as pdf:
        for page in pdf.pages:
            text += (page.extract_text() or "") + "\n"
    return text


def extract_text_from_docx(file: UploadFile) -> str:
    data = file.file.read()
    doc = DocxDocument(io.BytesIO(data))
    parts = [p.text for p in doc.paragraphs if p.text]
    for table in doc.tables:
        for row in table.rows:
            cells = [c.text.strip() for c in row.cells if c.text]
            if cells:
                parts.append(" | ".join(cells))
    return "\n".join(parts)


def extract_text_from_pptx(file: UploadFile) -> str:
    data = file.file.read()
    prs = Presentation(io.BytesIO(data))
    parts = []
    for idx, slide in enumerate(prs.slides, start=1):
        parts.append(f"\n--- Slide {idx} ---")
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    line = "".join(run.text for run in para.runs).strip()
                    if line:
                        parts.append(line)
    return "\n".join(parts)


def extract_text_from_txt(file: UploadFile) -> str:
    data = file.file.read()
    for enc in ("utf-8", "utf-16", "latin-1"):
        try:
            return data.decode(enc)
        except UnicodeDecodeError:
            continue
    return data.decode("utf-8", errors="ignore")


def extract_text_from_file(file: UploadFile) -> str:
    ext = _ext(file.filename)
    if ext == ".pdf":
        text = extract_text_from_pdf(file)
    elif ext == ".docx":
        text = extract_text_from_docx(file)
    elif ext == ".pptx":
        text = extract_text_from_pptx(file)
    elif ext == ".txt":
        text = extract_text_from_txt(file)
    elif ext in (".doc", ".ppt"):
        raise HTTPException(
            400,
            f"Legacy {ext} format is not supported. Please save as {ext}x and retry.",
        )
    else:
        raise HTTPException(
            400,
            f"Unsupported file type '{ext}'. Allowed: {', '.join(sorted(SUPPORTED_EXTENSIONS))}",
        )
    text = (text or "").strip()
    # Collapse runs of 3+ identical letters — some PDF/PPTX files extract with
    # each character duplicated (e.g. "TTTThhhhrrrree" -> "Three"). No English
    # word has a letter 3+ times in a row, so legit doubles are preserved.
    text = re.sub(r"([A-Za-z])\1{2,}", r"\1", text)
    # RAG handles full documents via chunking, so keep a generous cap
    # (protects memory on small/free hosting tiers).
    return text[:200000]


# ============================
# RAG CORE — CHUNK / EMBED / SEARCH
# ============================

VECTOR_INDEX_NAME = "vector_index"   # must match the Atlas Search index name
EMBED_MODEL = "jina-embeddings-v3"
EMBED_DIM = 1024                     # jina-embeddings-v3 default output size
TOP_K = 5                            # how many chunks to retrieve per question


def chunk_text(text: str, size: int = 1200, overlap: int = 200) -> List[str]:
    """Split text into overlapping, word-aware chunks."""
    text = (text or "").strip()
    if not text:
        return []

    words = text.split()
    chunks = []
    step = max(1, size - overlap)
    for start in range(0, len(words), step):
        chunk = " ".join(words[start:start + size])
        if chunk.strip():
            chunks.append(chunk)
    return chunks


def embed(texts: List[str], is_query: bool = False) -> List[List[float]]:
    """Get embeddings from Jina. is_query=True for the question, False for documents."""
    if not JINA_API_KEY:
        raise HTTPException(500, "JINA_API_KEY missing — embeddings not configured.")
    if not texts:
        return []

    payload = {
        "model": EMBED_MODEL,
        "task": "retrieval.query" if is_query else "retrieval.passage",
        "dimensions": EMBED_DIM,
        "input": [{"text": t} for t in texts],
    }
    try:
        resp = httpx.post(
            "https://api.jina.ai/v1/embeddings",
            headers={
                "Authorization": f"Bearer {JINA_API_KEY}",
                "Content-Type": "application/json",
            },
            json=payload,
            timeout=60,
        )
        resp.raise_for_status()
    except httpx.HTTPError as e:
        raise HTTPException(502, f"Embedding request failed: {e}")

    data = resp.json().get("data", [])
    # Keep the order Jina returns (it includes an "index" field)
    data.sort(key=lambda d: d.get("index", 0))
    return [d["embedding"] for d in data]


def vector_search(query_vector: List[float], document_id: str, user_id: str, k: int = TOP_K):
    """Find the most relevant chunks for a question, scoped to one user's document."""
    pipeline = [
        {
            "$vectorSearch": {
                "index": VECTOR_INDEX_NAME,
                "path": "embedding",
                "queryVector": query_vector,
                "numCandidates": max(100, k * 20),
                "limit": k,
                "filter": {"document_id": document_id, "user_id": user_id},
            }
        },
        {
            "$project": {
                "_id": 0,
                "text": 1,
                "chunk_index": 1,
                "score": {"$meta": "vectorSearchScore"},
            }
        },
    ]
    return list(chunks_collection.aggregate(pipeline))


def ask_groq(prompt: str, context: str = ""):
    if not groq_client:
        return "Groq API not configured."

    # 🧠 DETECT MIND MAP
    is_mind_map = "mind map" in prompt.lower()

    if is_mind_map:
        full_prompt = f"""
You are a strict mind map generator.

FOLLOW THESE RULES VERY STRICTLY:

1. Output ONLY tree structure
2. NEVER use bullet points (•)
3. NEVER write paragraphs
4. NEVER explain anything
5. Use ONLY these symbols: ├── └── │
6. Maintain proper indentation

EXAMPLE:

Machine Learning
├── Definition
│   ├── Learns from data
│   └── Improves with experience
├── Types
│   ├── Supervised
│   ├── Unsupervised
│   └── Reinforcement

NOW CONVERT THIS INTO SAME FORMAT:

{prompt}
"""
    else:
        full_prompt = f"""
You are a helpful AI assistant. Answer the question using ONLY the context.

STRICT OUTPUT RULES — follow exactly:
- Reply in PLAIN TEXT only. No Markdown syntax at all.
- Do NOT use tables, pipes (|), dashes as separators (---), or column layouts.
- Do NOT use **, __, ##, backticks, HTML tags, or LaTeX (no \\rightarrow, \\alpha, $...$ etc.).
- Write short paragraphs separated by a single blank line.
- For lists, put each item on its own line starting with "- " (dash + space).
- For sections, put a short plain-text title on its own line, then a blank line, then the content.
- Keep it concise, well organized, and easy to read.

Context:
{context[:3000]}

Question:
{prompt}

Answer (plain text only):
"""

    res = groq_client.chat.completions.create(
        model="openai/gpt-oss-120b",
        messages=[{"role": "user", "content": full_prompt}],
        temperature=0.2,
        max_tokens=500
    )

    return res.choices[0].message.content

# ============================
# ROUTES
# ============================

@app.get("/")
def root():
    return {"status": "running"}


@app.get("/health")
def health():
    return {
        "mongodb": "connected" if db is not None else "fallback",
        "groq": "ready" if groq_client else "missing"
    }


# ============================
# REGISTER
# ============================

@app.post("/register")
def register(user: UserRegister):
    if users_collection.find_one({"email": user.email}):
        raise HTTPException(400, "User exists")

    user_id = str(uuid.uuid4())

    users_collection.insert_one({
        "_id": user_id,
        "username": user.username,
        "email": user.email,
        "password": hash_password(user.password),
        "created_at": datetime.now()
    })

    return {"user_id": user_id}


# ============================
# LOGIN
# ============================

@app.post("/login")
def login(user: UserLogin):
    db_user = users_collection.find_one({"email": user.email})

    if not db_user or not verify_password(user.password, db_user["password"]):
        raise HTTPException(401, "Invalid credentials")

    return {
        "user_id": str(db_user["_id"]),
        "username": db_user["username"]
    }


# ============================
# UPLOAD DOCUMENT (PDF / DOCX / PPTX / TXT)
# ============================

async def _handle_upload(
    file: UploadFile,
    user_id: str,
    document_name: str,
):
    ext = _ext(file.filename)
    if ext not in SUPPORTED_EXTENSIONS and ext not in (".doc", ".ppt"):
        raise HTTPException(
            400,
            f"Unsupported file type '{ext}'. Allowed: {', '.join(sorted(SUPPORTED_EXTENSIONS))}",
        )

    text = extract_text_from_file(file)

    if not text:
        raise HTTPException(400, "Could not extract any text from the file.")

    doc_id = str(uuid.uuid4())

    # 1) Split the document into chunks
    chunks = chunk_text(text)
    if not chunks:
        raise HTTPException(400, "Could not split the document into chunks.")

    # 2) Embed every chunk via Jina
    vectors = embed(chunks, is_query=False)
    if len(vectors) != len(chunks):
        raise HTTPException(502, "Embedding count did not match chunk count.")

    # 3) Store one record per chunk (with its vector) for retrieval
    chunk_docs = [
        {
            "_id": str(uuid.uuid4()),
            "document_id": doc_id,
            "user_id": user_id,
            "chunk_index": i,
            "text": chunk,
            "embedding": vector,
        }
        for i, (chunk, vector) in enumerate(zip(chunks, vectors))
    ]
    chunks_collection.insert_many(chunk_docs)

    # Keep document metadata (content kept for the summary / reference)
    documents_collection.insert_one({
        "_id": doc_id,
        "user_id": user_id,
        "name": document_name,
        "file_type": ext.lstrip("."),
        "content": text,
        "num_chunks": len(chunks),
        "uploaded_at": datetime.now(),
    })

    # Brief summary from the start of the document (cheap, optional)
    summary = ask_groq("Summarize this document briefly", text[:3000])

    return {
        "document_id": doc_id,
        "document_name": document_name,
        "file_type": ext.lstrip("."),
        "num_chunks": len(chunks),
        "summary": summary[:200],
        "message": "Document uploaded successfully",
    }


@app.post("/upload/pdf")
async def upload_pdf(
    file: UploadFile = File(...),
    user_id: str = Form(...),
    document_name: str = Form("Untitled"),
):
    return await _handle_upload(file, user_id, document_name)


@app.post("/upload/document")
async def upload_document(
    file: UploadFile = File(...),
    user_id: str = Form(...),
    document_name: str = Form("Untitled"),
):
    return await _handle_upload(file, user_id, document_name)


# ============================
# CHAT
# ============================

@app.post("/chat/ai")
def chat(chat: ChatRequest):

    context = ""

    # If a document is selected → real RAG: embed the question, retrieve
    # the most relevant chunks, and build context only from those.
    if chat.document_id:
        doc = documents_collection.find_one({
            "_id": chat.document_id,
            "user_id": chat.user_id
        })

        if not doc:
            raise HTTPException(404, "Document not found")

        # 1) Embed the question
        query_vec = embed([chat.message], is_query=True)[0]

        # 2) Retrieve top-K relevant chunks via Atlas Vector Search
        try:
            hits = vector_search(query_vec, chat.document_id, chat.user_id)
        except Exception as e:
            # Most common cause: the Atlas vector index isn't created yet.
            raise HTTPException(
                500,
                f"Vector search failed (is the '{VECTOR_INDEX_NAME}' Atlas index created?): {e}",
            )

        # 3) Fallback for documents uploaded before RAG (no chunks yet)
        if not hits:
            context = doc.get("content", "")[:3000]
        else:
            context = "\n\n".join(h["text"] for h in hits)

    # Ask Groq using only the retrieved context
    answer = ask_groq(chat.message, context)

    responses_collection.insert_one({
        "_id": str(uuid.uuid4()),
        "user_id": chat.user_id,
        "document_id": chat.document_id,  # can be None
        "question": chat.message,
        "answer": answer,
        "timestamp": datetime.now()
    })

    return {
    "answer": answer,
    "source": context[:300]
}

# ============================
# ⭐ CHAT HISTORY (REQUIRED for React)
# ============================

@app.get("/chat/history/{user_id}")
def get_chat_history(user_id: str):

    chats = list(responses_collection.find({"user_id": user_id}))

    for chat in chats:
        chat["_id"] = str(chat["_id"])

    return {"chats": chats}



# ============================
# RUN
# ============================

# if __name__ == "__main__":
#     import uvicorn
#     uvicorn.run(app, host="0.0.0.0",port=8000, reload=True)
